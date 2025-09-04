import os
import zipfile
import requests
import json
import pymongo
from pptx import Presentation
import re

class OllamaPPTAnalyzer:
    def __init__(self, ollama_host="http://localhost:11434", model_name="llama2"):
        self.ollama_host = ollama_host
        self.model_name = model_name
        self.ollama_url = f"{ollama_host}/api/generate"
        print(f"[INIT] Initializing OllamaPPTAnalyzer...")
        print(f"[INIT] Ollama Host: {ollama_host}")
        print(f"[INIT] Model: {model_name}")
        
        self.client = pymongo.MongoClient('mongodb://localhost:27017/')
        self.db = self.client['internship-program']
        self.collection = self.db['ppt_analysis']
        print(f"[INIT] Connected to MongoDB: internship-program.ppt_analysis")

        if not self.test_ollama_connection():
            print("[ERROR] Ollama is not responding! Check if it's running properly.")
            exit(1)

    def test_ollama_connection(self):
        try:
            print(f"[INIT] Testing Ollama connection...")
            response = requests.get(f"{self.ollama_host}/api/version", timeout=5)
            if response.status_code == 200:
                print(f"[INIT] ✅ Ollama is responding")
                return True
            else:
                print(f"[INIT] ❌ Ollama returned status {response.status_code}")
                return False
        except Exception as e:
            print(f"[INIT] ❌ Cannot connect to Ollama: {e}")
            return False

    def clean_title(self, title):
        if not title:
            return "Unknown Title"
        # Remove common unwanted keywords
        unwanted_keywords = [
            r'2024\s*',
            r'ISG\s*',
            r'Hackathon\s*',
            r'PoC\s*',
            r'Template\s*',
            r'_.*$',  # Remove everything after underscore
        ]
        
        cleaned = title
        for pattern in unwanted_keywords:
            cleaned = re.sub(pattern, '', cleaned, flags=re.IGNORECASE)
        
        # Clean up extra spaces and punctuation
        cleaned = re.sub(r'\s+', ' ', cleaned).strip()
        cleaned = cleaned.strip('_-.,')
        
        return cleaned if cleaned else "Unknown Title"

    def analyze_ppt(self, ppt_path, max_retries=3):
        print(f"\n[PROCESSING] {os.path.basename(ppt_path)}")
        
        prs = Presentation(ppt_path)
        raw_title = None
        description = ""
        
        # Extract title from first slide
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    raw_title = shape.text_frame.text.strip().replace('\n', ' ').replace('\x0b', '')
                    break
        
        # Extract all text for context
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    description += shape.text_frame.text.strip().replace('\n', ' ').replace('\x0b', '') + " "
        
        # Clean the title
        clean_title = self.clean_title(raw_title)
        
        prompt = f"""
        You are analyzing a technical presentation. Based on the content below, provide a clear analysis.

        Raw Title: {raw_title}
        Content: {description}

        Provide a JSON response with these exact fields:
        {{
            "title": "A clear, concise title that describes what needs to be built/solved (remove years, hackathon, ISG keywords)",
            "description": "A detailed description explaining the problem statement and what exactly needs to be implemented or solved. Make it clear enough that a developer can understand the requirements.",
            "tech_stack": ["List of specific technologies, programming languages, frameworks, tools needed"],
            "skills_required": {{"skill_name": importance_score}}
        }}

        Requirements:
        - title: Clean, meaningful title without buzzwords
        - description: 2-3 sentences explaining the problem and solution approach
        - tech_stack: Specific technologies (e.g., "Python", "React", "MongoDB", "Docker")
        - skills_required: Skills with scores 1-100 based on importance

        Return only valid JSON, no extra text.
        """
        
        payload = {
            "model": self.model_name,
            "prompt": prompt,
            "stream": False,
            "format": "json"
        }
        
        # Test Ollama with a simple prompt first
        test_payload = {
            "model": self.model_name,
            "prompt": "Hello",
            "stream": False
        }

        try:
            print(f"  [TEST] Testing Ollama with simple prompt...")
            test_response = requests.post(self.ollama_url, json=test_payload, timeout=120)
            test_response.raise_for_status()
            print(f"  [TEST] ✅ Ollama is working")
        except Exception as e:
            print(f"  [TEST] ❌ Ollama test failed: {e}")
            return None

        for attempt in range(1, max_retries + 1):
            try:
                print(f"  [OLLAMA] Sending to Ollama (attempt {attempt})...")
                response = requests.post(self.ollama_url, json=payload, timeout=120)  # Changed from 120
                response.raise_for_status()
                
                result = response.json()
                generated_text = result.get('response', '')
                
                print(f"  [OLLAMA] Response received")
                analysis = json.loads(generated_text)
                analysis['ppt_name'] = os.path.basename(ppt_path)
                
                print(f"  [SUCCESS] Title: {analysis.get('title', 'N/A')}")
                print(f"  [SUCCESS] Tech Stack: {len(analysis.get('tech_stack', []))} items")
                print(f"  [SUCCESS] Skills: {len(analysis.get('skills_required', {}))} skills")
                
                return analysis
                
            except requests.exceptions.RequestException as e:
                print(f"  [ERROR] Ollama connection error: {e}")
                return None
            except json.JSONDecodeError as e:
                print(f"  [ERROR] JSON parsing failed (attempt {attempt}): {e}")
                if attempt < max_retries:
                    print(f"  [RETRY] Waiting 2 seconds...")
                    import time
                    time.sleep(2)
                    continue
                else:
                    print(f"  [ERROR] Max retries reached")
                    return None
            except Exception as e:
                print(f"  [ERROR] Unexpected error: {e}")
                return None

    def save_single_to_mongodb(self, analysis):
        try:
            result = self.collection.insert_one(analysis)
            print(f"  [MONGODB] ✅ Saved to database")
            return True
        except pymongo.errors.DuplicateKeyError:
            print(f"  [MONGODB] Already exists in database")
            return True
        except Exception as e:
            print(f"  [MONGODB] ❌ Error saving: {e}")
            return False

def main():
    print("="*60)
    print("[INFO] Starting PPT Analysis with Ollama")
    print("="*60)
    
    zip_path = r'd:\WorkSpace-Lenovo\projects\POC\internship_project\src\poc_ideas.zip'
    extract_dir = r'd:\WorkSpace-Lenovo\projects\POC\internship_project\src\poc_ideas_unzipped'
    
    if not os.path.exists(extract_dir):
        os.makedirs(extract_dir)
    
    # Unzip the file
    print(f"[SETUP] Extracting ZIP file...")
    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
    
    # Find all pptx files
    ppt_files = []
    for root, dirs, files in os.walk(extract_dir):
        for file in files:
            if file.lower().endswith('.pptx'):
                ppt_files.append(os.path.join(root, file))
    
    print(f"[SETUP] Found {len(ppt_files)} PPT files")
    
    analyzer = OllamaPPTAnalyzer(
        ollama_host="http://localhost:11434",
        model_name="llama2"
    )
    
    successful_analyses = 0
    failed_analyses = 0
    
    for i, ppt_path in enumerate(ppt_files, 1):
        print(f"\n[{i}/{len(ppt_files)}] Processing: {os.path.basename(ppt_path)}")
        
        analysis = analyzer.analyze_ppt(ppt_path)
        if analysis:
            analysis['ppt_id'] = i
            if analyzer.save_single_to_mongodb(analysis):
                successful_analyses += 1
            else:
                failed_analyses += 1
        else:
            print(f"  [FAILED] ❌ Analysis failed")
            failed_analyses += 1
    
    print(f"\n{'='*60}")
    print(f"[FINAL] Analysis Complete!")
    print(f"[FINAL] Successful: {successful_analyses}")
    print(f"[FINAL] Failed: {failed_analyses}")
    print(f"[FINAL] Total: {len(ppt_files)}")
    print("="*60)

if __name__ == "__main__":
    main()